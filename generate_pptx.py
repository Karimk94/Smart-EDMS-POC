from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
# Set to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Styling constants
FONT_NAME = 'Segoe UI'
TEXT_COLOR = RGBColor(0x24, 0x29, 0x2e) # GitHub dark gray
TITLE_COLOR = RGBColor(0x24, 0x29, 0x2e)
BLUE_LINK = RGBColor(0x03, 0x66, 0xd6)

def apply_font(run, size, bold=False, color=TEXT_COLOR):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def create_slide_with_header(title):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    apply_font(run, size=44, bold=True, color=TITLE_COLOR)
    
    # Divider line
    line = slide.shapes.add_shape(
        9, # msoShapeLine
        Inches(0.8), Inches(1.6), Inches(11.7), Inches(0)
    )
    line.line.color.rgb = RGBColor(0xea, 0xea, 0xef)
    line.line.width = Pt(2)
    
    return slide

def add_title_slide(title, subtitle):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    apply_font(run, size=54, bold=True, color=TITLE_COLOR)
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.3), Inches(1))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    apply_font(run2, size=32, bold=False, color=RGBColor(0x6a, 0x73, 0x7d))

def add_bullet_slide(title, text_lines):
    slide = create_slide_with_header(title)
    
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(text_lines):
        if line.strip() == "":
            p = tf.add_paragraph()
            p.space_before = Pt(14)
            continue
            
        p = tf.add_paragraph()
        p.space_before = Pt(14)
        
        is_bullet = line.startswith('•') or line.startswith('1.') or line.startswith('2.') or line.startswith('3.') or line.startswith('4.')
        is_sub = line.startswith('  -')
        
        if is_bullet:
            p.level = 0
        elif is_sub:
            p.level = 1
            line = line.replace('  -', '•')
        else:
            p.level = 0
            
        run = p.add_run()
        run.text = line
        apply_font(run, size=28, color=TEXT_COLOR)

def add_image_slide(title, text_before, image_path, bullet_lines):
    slide = create_slide_with_header(title)
    
    # Text Before
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text_before
    apply_font(run, size=24, color=TEXT_COLOR)
    
    # Image
    try:
        pic = slide.shapes.add_picture(image_path, Inches(1.6), Inches(2.3), width=Inches(10.1))
    except Exception as e:
        print(f"Could not add image {image_path}: {e}")
        
    # Bullets below
    if bullet_lines:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(bullet_lines):
            p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = line
            apply_font(run, size=20, color=TEXT_COLOR)

add_title_slide("Smart EDMS Photo AI Platform", "User Guide & System Overview")

add_bullet_slide("What is Smart EDMS?", [
    "Smart EDMS is an intelligent photo management system. It doesn't just store your photos; it understands them using advanced AI.",
    "",
    "• Organize: Folders, rename, and search.",
    "• Analyze: Automatically detects objects, faces, and text.",
    "• Recognize: Learn and recognize people's faces across your photo library."
])

add_image_slide("1. Navigating the Main Page", 
    "The main dashboard is where you view all your folders and photos.",
    r"C:\Users\okool_kaabdulwahed\.gemini\antigravity-ide\brain\5a6f7e81-5e52-43fe-8727-b6f678934306\main_page_v4_1780649385374.png",
    [
        "• Search Bar (Top): Type to instantly find folders or photos.",
        "• Language & Theme: Switch between English/Arabic and Light/Dark mode.",
        "• Create Folder / Upload: Quick action buttons in the top right."
    ]
)

add_image_slide("2. Uploading Photos", 
    "To add new photos, click the Upload button at the top right of the screen.",
    r"C:\Users\okool_kaabdulwahed\.gemini\antigravity-ide\brain\5a6f7e81-5e52-43fe-8727-b6f678934306\upload_modal_v4_1780649395188.png",
    [
        "• Drag & Drop: Simply drag photos from your computer into the dashed area.",
        "• Browse: Click 'Browse Photos' to select files manually.",
        "• Upload Button: Click the blue 'Upload Photos' button to start the process."
    ]
)

add_image_slide("3. Viewing and Analyzing a Photo", 
    "Click on any photo in your folder to open the Photo Details View.",
    r"C:\Users\okool_kaabdulwahed\.gemini\antigravity-ide\brain\5a6f7e81-5e52-43fe-8727-b6f678934306\photo_modal_v4_1780649414853.png",
    [
        "When a photo is uploaded, the AI automatically analyzes it in the background!"
    ]
)

add_bullet_slide("4. Understanding AI Analysis Results", [
    "In the Photo Details View (Right Panel), you will see the AI's findings:",
    "",
    "1. Caption: A human-like description of what is happening in the photo.",
    "2. Objects / Tags: Keywords of items detected (e.g., 'Car', 'Tree', 'Building').",
    "3. OCR Text: Any readable text found within the image (signs, documents, etc.).",
    "4. Detected Faces: Thumbnails of all faces found in the picture."
])

add_bullet_slide("5. Teaching the System to Recognize Faces", [
    "You can train the AI to recognize specific people!",
    "",
    "1. In the 'Detected Faces' section, find an unknown face.",
    "2. Click the 'Name & Save' button next to it.",
    "3. A popup will appear. Type the person's name or select an existing one.",
    "4. Click 'Save Face'.",
    "",
    "Next time you upload a photo with this person, the system will automatically recognize them!"
])

add_bullet_slide("Re-running Analysis", [
    "If you need the AI to take another look at a photo:",
    "",
    "• Click the Re-run Analysis dropdown button at the top right of the Photo Details View.",
    "• You can run everything, or select a specific service:",
    "  - OCR: Extract text again.",
    "  - Face Detection: Scan for faces again.",
    "  - Captions & Objects: Regenerate descriptions and tags."
])

add_bullet_slide("Summary", [
    "1. Upload photos securely.",
    "2. Let the AI analyze contents automatically.",
    "3. Search for text, objects, or people effortlessly.",
    "4. Name faces to build your custom recognition database."
])

prs.save('presentation_styled.pptx')
